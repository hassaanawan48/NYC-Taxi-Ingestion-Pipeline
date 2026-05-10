#!/usr/bin/env python3
"""Generate the project presentation (PPTX) for Big Data Analytics A2+A3.

Minimalistic, low-saturation theme. Slide structure:
  1. Title / course / group / contributions
  2. Dataset overview & rationale
  3. Data quality issues & cleaning
  4. Data analytics & visualization (multi-slide section)
  5. Business insights & conclusion
  6. Dataset & GitHub links
  Plus extra slides: pipeline architecture, warehouse star schema,
  pipeline optimizations, key numbers at a glance.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

PROJECT_DIR = Path(__file__).resolve().parent
OUTPUT_FILE = PROJECT_DIR / "NYC_Taxi_BDA_Presentation.pptx"

# Minimal palette: dark navy, charcoal, soft slate, single muted accent
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PANEL = RGBColor(0xF5, 0xF6, 0xF8)
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)        # deep navy for headings
COLOR_ACCENT = RGBColor(0xC9, 0xA2, 0x27)         # muted gold accent line
COLOR_TEXT = RGBColor(0x22, 0x2A, 0x35)           # almost-black charcoal
COLOR_MUTED = RGBColor(0x6B, 0x73, 0x80)          # subtle gray for captions
COLOR_RULE = RGBColor(0xD7, 0xDC, 0xE3)           # soft divider gray

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_solid_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT,
    align=PP_ALIGN.LEFT,
    font: str = "Calibri",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    items,
    *,
    size: int = 16,
    color: RGBColor = COLOR_TEXT,
    bullet: str = "•  ",
    line_spacing: float = 1.25,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for idx, item in enumerate(items):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = line_spacing
        if isinstance(item, tuple):
            head, body = item
            r1 = para.add_run()
            r1.text = bullet + head
            r1.font.name = "Calibri"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = COLOR_PRIMARY
            r2 = para.add_run()
            r2.text = " — " + body
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.bold = False
            r2.font.color.rgb = color
        else:
            run = para.add_run()
            run.text = bullet + str(item)
            run.font.name = "Calibri"
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return tb


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_solid_fill(shape, color)
    return shape


def _slide_chrome(slide, slide_no: int, total: int, section: str | None = None):
    """Apply the minimal frame: white background, accent rule, footer."""
    bg = _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLOR_BG)
    bg.shadow.inherit = False

    # Slim left accent bar
    _add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, COLOR_PRIMARY)

    # Footer rule
    _add_rect(slide, Inches(0.5), Inches(7.05), Inches(12.4), Emu(9525), COLOR_RULE)

    # Footer text
    if section:
        _add_text(
            slide,
            Inches(0.5),
            Inches(7.12),
            Inches(8.0),
            Inches(0.3),
            section,
            size=10,
            color=COLOR_MUTED,
        )
    _add_text(
        slide,
        Inches(11.5),
        Inches(7.12),
        Inches(1.5),
        Inches(0.3),
        f"{slide_no} / {total}",
        size=10,
        color=COLOR_MUTED,
        align=PP_ALIGN.RIGHT,
    )
    _add_text(
        slide,
        Inches(0.5),
        Inches(7.32),
        Inches(8.0),
        Inches(0.25),
        "NYC Yellow Taxi Pipeline · Group 17 · BSCS-13A",
        size=9,
        color=COLOR_MUTED,
    )


def _slide_header(slide, kicker: str, title: str):
    _add_text(
        slide,
        Inches(0.55),
        Inches(0.45),
        Inches(11),
        Inches(0.35),
        kicker.upper(),
        size=11,
        bold=True,
        color=COLOR_ACCENT,
    )
    _add_text(
        slide,
        Inches(0.55),
        Inches(0.78),
        Inches(12),
        Inches(0.85),
        title,
        size=32,
        bold=True,
        color=COLOR_PRIMARY,
    )
    _add_rect(slide, Inches(0.55), Inches(1.55), Inches(0.6), Emu(28575), COLOR_ACCENT)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    total_slides = 11

    # ================== SLIDE 1: TITLE ==================
    s = prs.slides.add_slide(blank_layout)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_BG)
    _add_rect(s, 0, 0, Inches(0.12), SLIDE_H, COLOR_PRIMARY)
    _add_rect(s, Inches(0.6), Inches(0.6), Inches(2.6), Emu(28575), COLOR_ACCENT)

    _add_text(
        s,
        Inches(0.6),
        Inches(0.75),
        Inches(12),
        Inches(0.4),
        "BIG DATA ANALYTICS  ·  ASSIGNMENTS 2 & 3",
        size=12,
        bold=True,
        color=COLOR_ACCENT,
    )
    _add_text(
        s,
        Inches(0.6),
        Inches(1.2),
        Inches(12),
        Inches(1.2),
        "NYC Yellow Taxi",
        size=54,
        bold=True,
        color=COLOR_PRIMARY,
    )
    _add_text(
        s,
        Inches(0.6),
        Inches(2.0),
        Inches(12),
        Inches(0.9),
        "Ingestion, ETL & Analytics on HDFS + Spark",
        size=28,
        color=COLOR_TEXT,
    )

    # Course / class panel
    _add_rect(s, Inches(0.6), Inches(3.2), Inches(5.9), Inches(2.4), COLOR_PANEL)
    _add_text(
        s, Inches(0.85), Inches(3.35), Inches(5.5), Inches(0.35),
        "COURSE & CLASS", size=11, bold=True, color=COLOR_ACCENT,
    )
    _add_paragraphs(
        s, Inches(0.85), Inches(3.7), Inches(5.5), Inches(2.0),
        [
            ("Course", "Big Data Analytics"),
            ("Class", "BSCS-13A"),
            ("Course Group", "Group 3 (BDA, BSCS-13A)"),
            ("Project Group #", "17"),
            ("Institution", "NUST SEECS"),
        ],
        size=15,
        bullet="",
    )

    # Members panel
    _add_rect(s, Inches(6.85), Inches(3.2), Inches(5.9), Inches(2.4), COLOR_PANEL)
    _add_text(
        s, Inches(7.1), Inches(3.35), Inches(5.5), Inches(0.35),
        "GROUP 17  ·  CONTRIBUTIONS", size=11, bold=True, color=COLOR_ACCENT,
    )
    _add_paragraphs(
        s, Inches(7.1), Inches(3.7), Inches(5.5), Inches(2.0),
        [
            ("Malik Muhammad Hassaan", "CMS 460233 · Data ingestion, validation, profiling"),
            ("Muhammad Ahmed", "CMS 456679 · ETL pipeline, business-question SQL, optimizations"),
        ],
        size=14,
        bullet="",
        line_spacing=1.4,
    )

    # Footer note
    _add_rect(s, Inches(0.6), Inches(6.55), Inches(12.2), Emu(28575), COLOR_RULE)
    _add_text(
        s, Inches(0.6), Inches(6.65), Inches(12.2), Inches(0.4),
        "End-to-end pipeline · 2.96M trips · HDFS warehouse · Spark SQL analytics · automated reporting",
        size=12, color=COLOR_MUTED, align=PP_ALIGN.LEFT,
    )

    # ================== SLIDE 2: DATASET OVERVIEW ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 2, total_slides, "Dataset")
    _slide_header(s, "Slide 2", "Dataset Overview & Rationale")

    _add_paragraphs(
        s, Inches(0.55), Inches(1.95), Inches(7.2), Inches(4.8),
        [
            ("Source", "NYC TLC Yellow Taxi Trip Records"),
            ("File", "yellow_tripdata_2024-01.parquet"),
            ("Period", "January 2024 (one full month)"),
            ("Volume", "2,964,624 raw trip records"),
            ("Schema", "19 columns: datetime, numeric, categorical, location IDs"),
            ("Format", "Parquet (binary columnar, compressed)"),
            ("Provider", "nyc.gov / TLC Trip Record Data"),
        ],
        size=15,
        line_spacing=1.3,
    )

    # Right-side rationale panel
    _add_rect(s, Inches(8.0), Inches(1.95), Inches(4.85), Inches(4.8), COLOR_PANEL)
    _add_text(
        s, Inches(8.25), Inches(2.05), Inches(4.5), Inches(0.4),
        "WHY THIS DATASET", size=11, bold=True, color=COLOR_ACCENT,
    )
    _add_paragraphs(
        s, Inches(8.25), Inches(2.45), Inches(4.5), Inches(4.0),
        [
            "Volume > 500K rows justifies HDFS + Spark",
            "Each row is a trip — a natural fact table",
            "Clear dimensions: time, location, payment, vendor",
            "Contains realistic data quality issues for ETL",
            "Public, well-documented, columnar format",
        ],
        size=14,
        line_spacing=1.3,
    )

    # ================== SLIDE 3: DATA QUALITY ISSUES ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 3, total_slides, "Data Quality")
    _slide_header(s, "Slide 3", "Data Quality Issues & Cleaning")

    _add_text(
        s, Inches(0.55), Inches(1.85), Inches(12), Inches(0.4),
        "Identified during profiling (profile_data.py); resolved during ETL (etl.py).",
        size=13, color=COLOR_MUTED,
    )

    issues_left = [
        ("Negative fare amounts", "impute with distance-bin median fare"),
        ("Zero passenger count", "set minimum to 1 passenger"),
        ("Non-integer passenger count", "round to nearest integer"),
        ("Duplicate trip rows", "dropDuplicates (keep first)"),
        ("Trip distance outliers", "cap at 99th percentile"),
    ]
    issues_right = [
        ("Tip greater than fare", "cap tip at cleaned fare amount"),
        ("Missing numeric values", "median imputation per column"),
        ("Missing categorical codes", "fallback to 0 / mode"),
        ("Datetime inconsistencies", "standardize to timestamp; derive duration"),
        ("Future pickup timestamps", "replace with hour/weekday median"),
    ]

    _add_rect(s, Inches(0.55), Inches(2.4), Inches(6.05), Inches(4.4), COLOR_PANEL)
    _add_text(s, Inches(0.8), Inches(2.5), Inches(5.6), Inches(0.4),
              "ISSUE  →  ACTION", size=11, bold=True, color=COLOR_ACCENT)
    _add_paragraphs(
        s, Inches(0.8), Inches(2.9), Inches(5.6), Inches(3.8),
        issues_left, size=13, line_spacing=1.35,
    )

    _add_rect(s, Inches(6.75), Inches(2.4), Inches(6.05), Inches(4.4), COLOR_PANEL)
    _add_text(s, Inches(7.0), Inches(2.5), Inches(5.6), Inches(0.4),
              "ISSUE  →  ACTION", size=11, bold=True, color=COLOR_ACCENT)
    _add_paragraphs(
        s, Inches(7.0), Inches(2.9), Inches(5.6), Inches(3.8),
        issues_right, size=13, line_spacing=1.35,
    )

    # ================== SLIDE 4: PIPELINE ARCHITECTURE (extra) ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 4, total_slides, "Pipeline")
    _slide_header(s, "Bonus", "Pipeline Architecture")

    stages = [
        ("Ingest", "ingest.py", "Download · validate · upload to HDFS"),
        ("Profile", "profile_data.py", "Schema · missing · stats · QA"),
        ("ETL", "etl.py (Spark)", "Clean · transform · model · partition"),
        ("Analytics", "analytics.py (Spark SQL)", "5 business questions · charts"),
        ("Report", "build_final_report.py", "Bundle CSVs + charts → PDF"),
    ]
    box_w = Inches(2.36)
    box_h = Inches(3.2)
    gap = Inches(0.18)
    start_left = Inches(0.55)
    top = Inches(2.1)

    for i, (title, sub, body) in enumerate(stages):
        left = Emu(int(start_left) + i * (int(box_w) + int(gap)))
        _add_rect(s, left, top, box_w, box_h, COLOR_PANEL)
        _add_rect(s, left, top, box_w, Inches(0.06), COLOR_ACCENT)
        _add_text(s, Emu(int(left) + int(Inches(0.2))), Emu(int(top) + int(Inches(0.25))),
                  Inches(2.0), Inches(0.5), title,
                  size=20, bold=True, color=COLOR_PRIMARY)
        _add_text(s, Emu(int(left) + int(Inches(0.2))), Emu(int(top) + int(Inches(0.85))),
                  Inches(2.0), Inches(0.4), sub,
                  size=11, bold=True, color=COLOR_ACCENT)
        _add_text(s, Emu(int(left) + int(Inches(0.2))), Emu(int(top) + int(Inches(1.25))),
                  Inches(2.0), Inches(1.8), body,
                  size=12, color=COLOR_TEXT)

    _add_text(
        s, Inches(0.55), Inches(5.6), Inches(12), Inches(0.4),
        "Storage layer: HDFS  /warehouse/raw/...  →  /warehouse/processed/...",
        size=14, color=COLOR_MUTED,
    )
    _add_text(
        s, Inches(0.55), Inches(6.0), Inches(12), Inches(0.4),
        "Compute layer: Apache Spark 3.5.8 (PySpark, Spark SQL, broadcast joins, partitioning, caching)",
        size=14, color=COLOR_MUTED,
    )

    # ================== SLIDE 5: STAR SCHEMA (extra) ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 5, total_slides, "Warehouse Model")
    _slide_header(s, "Bonus", "Dimensional Model (Star Schema)")

    # Center fact box
    fact_left = Inches(5.0)
    fact_top = Inches(3.0)
    fact_w = Inches(3.4)
    fact_h = Inches(1.8)
    _add_rect(s, fact_left, fact_top, fact_w, fact_h, COLOR_PRIMARY)
    _add_text(s, fact_left, Emu(int(fact_top) + int(Inches(0.18))), fact_w, Inches(0.5),
              "fact_trips", size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
              align=PP_ALIGN.CENTER)
    _add_text(s, fact_left, Emu(int(fact_top) + int(Inches(0.7))), fact_w, Inches(0.5),
              "2,964,624 rows", size=14, color=RGBColor(0xFF, 0xFF, 0xFF),
              align=PP_ALIGN.CENTER)
    _add_text(s, fact_left, Emu(int(fact_top) + int(Inches(1.1))), fact_w, Inches(0.5),
              "trip distance · fare · tip · revenue · IDs",
              size=11, color=COLOR_PANEL, align=PP_ALIGN.CENTER)

    dims = [
        ("dim_datetime", "749 rows", Inches(0.6), Inches(1.95)),
        ("dim_payment", "5 rows", Inches(9.6), Inches(1.95)),
        ("dim_trip_category", "3 rows", Inches(5.0), Inches(5.6)),
    ]
    for name, sub, left, top in dims:
        _add_rect(s, left, top, Inches(3.0), Inches(1.3), COLOR_PANEL)
        _add_rect(s, left, top, Inches(3.0), Inches(0.06), COLOR_ACCENT)
        _add_text(s, left, Emu(int(top) + int(Inches(0.2))), Inches(3.0), Inches(0.5),
                  name, size=18, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
        _add_text(s, left, Emu(int(top) + int(Inches(0.75))), Inches(3.0), Inches(0.5),
                  sub, size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    _add_text(
        s, Inches(0.55), Inches(6.55), Inches(12), Inches(0.35),
        "fact_trips joins to each dimension via natural keys (datetime_key, payment_type, distance_bucket).",
        size=12, color=COLOR_MUTED,
    )

    # ================== SLIDE 6: ANALYTICS — BUSINESS QUESTIONS (header) ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 6, total_slides, "Analytics")
    _slide_header(s, "Slide 4", "Data Analytics & Visualization — 5 Business Questions")

    bq = [
        "Q1.  How do trip volume and average fare vary by hour and day of week?",
        "Q2.  What is the correlation between trip distance and fare amount?",
        "Q3.  How does payment type preference change by time and location?",
        "Q4.  How does average tip differ by payment type, distance, or vendor?",
        "Q5.  Which pickup/dropoff location pairs have highest volume and fare?",
    ]
    _add_paragraphs(
        s, Inches(0.55), Inches(2.0), Inches(12.2), Inches(4.6),
        bq,
        size=18, line_spacing=1.5, bullet="",
        color=COLOR_TEXT,
    )

    _add_text(
        s, Inches(0.55), Inches(6.5), Inches(12), Inches(0.4),
        "Source: Dataset_Justification.pdf  ·  Implemented as Spark SQL queries in analytics.py.",
        size=12, color=COLOR_MUTED,
    )

    # ================== SLIDE 7: ANALYTICS — Q1 + Q2 (chart-driven) ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 7, total_slides, "Analytics")
    _slide_header(s, "Findings", "Q1 — Hourly Demand   ·   Q2 — Distance vs Fare")

    line_chart = PROJECT_DIR / "outputs" / "chart_line_trend.png"
    if line_chart.exists():
        s.shapes.add_picture(str(line_chart), Inches(0.55), Inches(2.0),
                             width=Inches(6.2), height=Inches(3.7))
    scatter_chart = PROJECT_DIR / "outputs" / "chart_scatter_revenue_duration.png"
    if scatter_chart.exists():
        s.shapes.add_picture(str(scatter_chart), Inches(7.05), Inches(2.0),
                             width=Inches(5.8), height=Inches(3.7))

    _add_paragraphs(
        s, Inches(0.55), Inches(5.85), Inches(6.0), Inches(1.3),
        [
            "Demand peaks 5–7 PM on weekdays.",
            "Fri/Sat sustain volume late at night.",
            "Lowest volume in the 3–5 AM window.",
        ],
        size=12, line_spacing=1.25,
    )
    _add_paragraphs(
        s, Inches(7.05), Inches(5.85), Inches(6.0), Inches(1.3),
        [
            "Strong positive distance-vs-fare relationship.",
            "Within-bin correlation low (per-bin variance).",
            "Confirms distance as primary fare driver.",
        ],
        size=12, line_spacing=1.25,
    )

    # ================== SLIDE 8: ANALYTICS — Q3, Q4, Q5 ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 8, total_slides, "Analytics")
    _slide_header(s, "Findings", "Q3 — Payment   ·   Q4 — Tipping   ·   Q5 — Top Pairs")

    bar_chart = PROJECT_DIR / "outputs" / "chart_bar_hour_revenue.png"
    if bar_chart.exists():
        s.shapes.add_picture(str(bar_chart), Inches(0.55), Inches(2.0),
                             width=Inches(6.2), height=Inches(3.6))

    _add_rect(s, Inches(7.05), Inches(2.0), Inches(5.8), Inches(3.6), COLOR_PANEL)
    _add_text(s, Inches(7.25), Inches(2.1), Inches(5.4), Inches(0.4),
              "KEY NUMBERS", size=11, bold=True, color=COLOR_ACCENT)
    _add_paragraphs(
        s, Inches(7.25), Inches(2.5), Inches(5.4), Inches(3.0),
        [
            ("Q3 Payment", "Credit card ≈ 80–84% in core Manhattan zones"),
            ("Q4 Tip (CC long)", "$11.58 avg · 21.2% of fare"),
            ("Q4 Tip (CC short)", "$2.71 avg · 27.9% of fare"),
            ("Q4 Tip (Cash)", "≈ 0 (No-charge / Dispute → ~$0.05)"),
            ("Q5 Top pair", "237 → 236  ·  21,883 trips"),
            ("Q5 Top revenue pair", "132 → 132  ·  $246,050 (JFK)"),
        ],
        size=13, line_spacing=1.3,
    )

    _add_paragraphs(
        s, Inches(0.55), Inches(5.7), Inches(12.2), Inches(1.3),
        [
            "Credit card dominates payment mix at high-volume zones during peak hours.",
            "Tipping rate (% of fare) is highest on short trips; tip dollar amount highest on long trips.",
            "A handful of pickup/dropoff pairs (Manhattan core + JFK) concentrate disproportionate revenue.",
        ],
        size=12, line_spacing=1.25, color=COLOR_TEXT,
    )

    # ================== SLIDE 9: PIPELINE OPTIMIZATIONS (extra) ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 9, total_slides, "Engineering")
    _slide_header(s, "Bonus", "Pipeline Optimizations")

    opts = [
        ("Broadcast joins", "Small dim tables broadcast to executors → eliminates shuffle on fact↔dim joins."),
        ("Caching reused DataFrame", "Cleaned DF cached before deriving fact + dimensions to avoid recomputation."),
        ("Partitioned Parquet writes", "fact_trips & dim_datetime partitioned by pickup_year, pickup_month."),
        ("percentile_approx", "Fast quantile estimation for distance capping & median imputation."),
        ("Dynamic partition overwrite", "spark.sql.sources.partitionOverwriteMode=dynamic for incremental loads."),
        ("Tuned shuffle partitions", "spark.sql.shuffle.partitions=16 for the single-month workload."),
    ]
    _add_paragraphs(
        s, Inches(0.55), Inches(1.95), Inches(12.2), Inches(4.8),
        opts, size=15, line_spacing=1.35,
    )

    # ================== SLIDE 10: BUSINESS INSIGHTS & CONCLUSION ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 10, total_slides, "Insights")
    _slide_header(s, "Slide 5", "Business Insights & Conclusion")

    insights = [
        ("Demand is highly time-bound", "Peak 5–7 PM and pre-dawn lulls suggest dynamic supply / surge windows."),
        ("Distance is the dominant fare driver", "Strong positive trend supports predictable, distance-based pricing."),
        ("NYC taxis are effectively cashless", "Credit card holds 80%+ share in the busiest zones — invest in card UX."),
        ("Tip behaviour is payment-driven", "Card riders tip materially; cash / no-charge segments tip near zero."),
        ("Revenue concentrates in a handful of OD pairs", "Manhattan core (161/236/237) and JFK (132) drive disproportionate revenue."),
        ("Data quality discipline pays off", "ETL preserved all 2.96M trips while neutralizing nulls, outliers and impossible values."),
    ]
    _add_paragraphs(
        s, Inches(0.55), Inches(1.95), Inches(12.2), Inches(4.9),
        insights, size=14, line_spacing=1.35,
    )

    # ================== SLIDE 11: LINKS ==================
    s = prs.slides.add_slide(blank_layout)
    _slide_chrome(s, 11, total_slides, "References")
    _slide_header(s, "Slide 6", "Dataset & GitHub Repository")

    _add_rect(s, Inches(0.55), Inches(2.0), Inches(12.2), Inches(2.0), COLOR_PANEL)
    _add_text(s, Inches(0.85), Inches(2.15), Inches(11.5), Inches(0.4),
              "DATASET", size=11, bold=True, color=COLOR_ACCENT)
    _add_text(s, Inches(0.85), Inches(2.55), Inches(11.5), Inches(0.5),
              "NYC TLC Yellow Taxi Trip Records (Jan 2024)",
              size=20, bold=True, color=COLOR_PRIMARY)
    _add_text(s, Inches(0.85), Inches(3.1), Inches(11.5), Inches(0.5),
              "https://www.nyc.gov/site/tlc/about/tlc-trip-record-data.page",
              size=14, color=COLOR_TEXT)
    _add_text(s, Inches(0.85), Inches(3.45), Inches(11.5), Inches(0.5),
              "Direct file: yellow_tripdata_2024-01.parquet",
              size=12, color=COLOR_MUTED)

    _add_rect(s, Inches(0.55), Inches(4.2), Inches(12.2), Inches(2.0), COLOR_PANEL)
    _add_text(s, Inches(0.85), Inches(4.35), Inches(11.5), Inches(0.4),
              "GITHUB REPOSITORY", size=11, bold=True, color=COLOR_ACCENT)
    _add_text(s, Inches(0.85), Inches(4.75), Inches(11.5), Inches(0.5),
              "NYC-Taxi-Ingestion-Pipeline",
              size=20, bold=True, color=COLOR_PRIMARY)
    _add_text(s, Inches(0.85), Inches(5.3), Inches(11.5), Inches(0.5),
              "https://github.com/hassaanawan48/NYC-Taxi-Ingestion-Pipeline",
              size=14, color=COLOR_TEXT)
    _add_text(s, Inches(0.85), Inches(5.65), Inches(11.5), Inches(0.5),
              "Includes ingest.py, profile_data.py, etl.py, analytics.py, build_final_report.py and PDFs.",
              size=12, color=COLOR_MUTED)

    _add_text(
        s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.4),
        "Thank you.  ·  Questions are welcome.",
        size=18, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER,
    )

    prs.save(OUTPUT_FILE)
    return OUTPUT_FILE


if __name__ == "__main__":
    out = build_presentation()
    print(f"Presentation saved to: {out}")
